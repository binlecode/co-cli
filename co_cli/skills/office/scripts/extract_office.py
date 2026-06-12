"""Extract a local Office file (.docx/.pptx/.xlsx) to markdown for the `office` skill.

Reached only as the ``co-extract-office`` console entry point through ``shell_exec`` —
never imported into the agent process. Dispatches by extension to a format-specific,
ML-free backend (``mammoth`` for docx, ``python-pptx`` for pptx, ``openpyxl`` for xlsx)
and emits markdown with format-appropriate structure markers: ``## Slide N`` for pptx,
``## Sheet <name>`` (a markdown table) for xlsx, and the document's own headings for docx
(only where Word heading styles exist — a directly-formatted docx degrades to flat prose,
with no synthetic anchors). On any error it writes a single plain line to stderr and exits
non-zero (no traceback).

An encrypted/password-protected OOXML file is an OLE/CFB compound document, not a Zip, so
every backend would fail it with a generic "not a zip" error indistinguishable from genuine
corruption. To give the password case its own message, the leading magic bytes are sniffed
up front (CFB header -> encrypted; a Zip header that still fails to parse -> corrupt).

Output goes through ``sys.stdout.write``/``sys.stderr.write`` (not ``print``) because this
module lives under ``co_cli/`` and ruff T20 forbids ``print`` there.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

SUPPORTED_SUFFIXES = (".docx", ".pptx", ".xlsx")

# An OLE/CFB compound document (legacy binary Office, or an encrypted OOXML wrapper).
CFB_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
# A Zip local-file header — the container every born-OOXML file actually is.
ZIP_MAGIC = b"PK\x03\x04"

# Cap rows rendered per xlsx sheet. openpyxl would otherwise serialize every cell, so a
# large sheet could emit an unbounded markdown table that swamps the agent's context. When
# a sheet exceeds the cap the script emits an explicit truncation line, so a capped read
# never looks complete (mirrors extract_pdf.py's total_pages notice).
DEFAULT_MAX_ROWS = 1000


def _fail(message: str) -> int:
    """Write a one-line message to stderr and return a non-zero exit code."""
    sys.stderr.write(message + "\n")
    return 1


def _classify_container(path: Path, raw: str) -> int:
    """Return 0 if ``path`` looks like a parseable OOXML Zip, else non-zero via _fail.

    Sniffs the leading magic bytes so the encrypted case (a CFB wrapper) gets its own
    message instead of surfacing as a generic backend parse error.
    """
    try:
        with path.open("rb") as handle:
            header = handle.read(8)
    except OSError:
        return _fail(f"Could not open Office file (corrupt or unreadable): {raw}")
    if header.startswith(CFB_MAGIC):
        return _fail(f"Office file is password-protected or encrypted: {raw}")
    if not header.startswith(ZIP_MAGIC):
        return _fail(f"Could not open Office file (corrupt or unreadable): {raw}")
    return 0


def _extract_docx(path: Path) -> str:
    """Convert a docx to markdown via mammoth (headings only where heading styles exist)."""
    import mammoth

    with path.open("rb") as handle:
        result = mammoth.convert_to_markdown(handle)
    return result.value.strip()


def _shape_text(shape) -> list[str]:
    """Collect text lines from one pptx shape (text frame and/or table)."""
    lines: list[str] = []
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            lines.append(text)
    if shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))
    return lines


def _extract_pptx(path: Path) -> str:
    """Render a pptx to markdown with a ``## Slide N`` (1-based) marker per slide."""
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            lines.extend(_shape_text(shape))
        body = "\n".join(lines).strip()
        parts.append(f"## Slide {index}\n\n{body}" if body else f"## Slide {index}")
    return "\n\n".join(parts)


def _cell_markdown(value) -> str:
    """Render one cell value for a markdown table (empty for blanks, pipes escaped)."""
    if value is None:
        return ""
    return str(value).replace("|", "\\|")


def _sheet_table(rows: list[list], total_rows: int, max_rows: int) -> str:
    """Build a markdown table from the rendered rows, appending a truncation line if capped.

    The first rendered row is the header. ``total_rows`` is the sheet's full row count
    before the cap, so the truncation notice reports the true denominator.
    """
    if not rows:
        return "(empty sheet)"
    header = rows[0]
    width = max(len(row) for row in rows)
    header = header + [""] * (width - len(header))
    lines = [
        "| " + " | ".join(_cell_markdown(cell) for cell in header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in rows[1:]:
        padded = list(row) + [None] * (width - len(row))
        lines.append("| " + " | ".join(_cell_markdown(cell) for cell in padded) + " |")
    table = "\n".join(lines)
    if total_rows > max_rows:
        table += f"\n\n[truncated: showing rows 1–{max_rows} of {total_rows}]"
    return table


def _extract_xlsx(path: Path, max_rows: int) -> str:
    """Render an xlsx to markdown with a ``## Sheet <name>`` table per worksheet."""
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for worksheet in workbook.worksheets:
            rendered: list[list] = []
            total_rows = 0
            for row in worksheet.iter_rows(values_only=True):
                total_rows += 1
                if len(rendered) < max_rows:
                    rendered.append(list(row))
            table = _sheet_table(rendered, total_rows, max_rows)
            parts.append(f"## Sheet {worksheet.title}\n\n{table}")
        return "\n\n".join(parts)
    finally:
        workbook.close()


def _extract(path: Path, suffix: str, max_rows: int) -> str:
    """Dispatch to the format-specific backend for a validated OOXML file."""
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    return _extract_xlsx(path, max_rows)


def main() -> int:
    """Console entry point (``co-extract-office``); returns the process exit code."""
    parser = argparse.ArgumentParser(
        prog="co-extract-office",
        description="Extract a local .docx/.pptx/.xlsx file to markdown.",
    )
    parser.add_argument("path", help="path to a local .docx, .pptx, or .xlsx file")
    parser.add_argument(
        "--max-rows",
        type=int,
        default=DEFAULT_MAX_ROWS,
        help=f"xlsx: cap rows rendered per sheet (default {DEFAULT_MAX_ROWS})",
    )
    args = parser.parse_args()

    path = Path(args.path)
    if not path.exists():
        return _fail(f"File not found: {args.path}")

    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        if suffix == ".pdf":
            return _fail(f"Not an Office file — use the documents skill for PDF: {args.path}")
        return _fail(
            f"Unsupported file type for office extraction "
            f"(expected .docx/.pptx/.xlsx): {args.path}"
        )

    invalid = _classify_container(path, args.path)
    if invalid:
        return invalid

    try:
        markdown = _extract(path, suffix, args.max_rows)
    except Exception:
        return _fail(f"Could not open Office file (corrupt or unreadable): {args.path}")

    sys.stdout.write(markdown + "\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
